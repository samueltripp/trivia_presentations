from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches

from pptx.enum.action import PP_ACTION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# method to create hyperlinks
def Hyperlink( run_object, source_slide, destination_slide ):
    rId = source_slide.part.relate_to(destination_slide.part, RT.SLIDE)
    rPr = run_object._r.get_or_add_rPr()
    hlinkClick = rPr.add_hlinkClick(rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


# name needs to be updated
name = "fall"
data = open(name + ".txt","r")
prs = Presentation("blank.pptx")


# Creating Title Slides, rules slides
title_layout = prs.slide_layouts[0]
title = prs.slides.add_slide(title_layout)
for i in range(2):
    par = title.shapes[i].text_frame.add_paragraph()
    par.font.color.rgb = RGBColor(0,0,0)
    par.text = data.readline()

# Creating pre-question slides
layout = prs.slide_layouts[1]
for i in range(2):
    # if there are no rules, creates nothing
    num_rules = int(data.readline().rstrip())
    if num_rules != 0:
        rules = prs.slides.add_slide(layout)
        p = rules.placeholders[0].text_frame.add_paragraph()
        p.text = "Rules"
        p.font.size = Pt(44)

        first_rule = rules.shapes[1].text_frame.paragraphs[0].add_run()
        first_rule.text = data.readline().rstrip()
        first_rule.font.size = Pt(28)

        for i in range(1, num_rules):
            tf = rules.shapes[1].text_frame
            r = tf.add_paragraph()
            r.text = data.readline().rstrip()
            r.font.size = Pt(28)

# read in the rest of the contest to deal with later, and find the category markers
content = data.readlines()
content = [line.rstrip() for line in content]
cat_index = []
for i in range(len(content)):
    if content[i].startswith("Category:"):
        cat_index.append(i)
        content[i] = content[i][9:]

#category title slide
cat_title_slide = len(list(prs.slides))
categories = prs.slides.add_slide(layout)
p = categories.placeholders[0].text_frame.add_paragraph()
p.text = "Today's Categories"
p.font.size = Pt(44)

# want to store the slide objects for categories, and the category indices fo musics and pictures
category_slides = []
music_cats=[]
picture_cats=[]

# create question slides
for i in range(len(cat_index)):
    music = False
    picture = False
    if content[cat_index[i]].endswith(":Music"):
        music = True
        music_cats.append(i)
        content[cat_index[i]] = content[cat_index[i]][:-6]
    if content[cat_index[i]].endswith(":Picture"):
        picture = True
        picture_cats.append(i)
        print(content[cat_index[i]])
        content[cat_index[i]] = content[cat_index[i]][:-8]

    # add category to "Today's Categories" slide
    if i == 0:
        p = categories.shapes[1].text_frame.paragraphs[0].add_run()
        p.text = content[cat_index[i]]+" "
        p.font.size = Pt(28)
    else:
        p = categories.shapes[1].text_frame.add_paragraph()
        p.text = content[cat_index[i]]+" "
        p.font.size = Pt(28)

    # make the category slide
    front = prs.slides.add_slide(layout)
    category_slides.append(front)
    p = front.placeholders[0].text_frame.add_paragraph()
    p.text = content[cat_index[i]]
    p.font.size = Pt(44)
    desc = front.shapes[1].text_frame.paragraphs[0].add_run()
    desc.text = content[cat_index[i]+1]
    desc.font.size = Pt(28)

    # link to the Today's Categories slide
    r = categories.shapes[1].text_frame.paragraphs[i].add_run()
    r.text = "Questions"
    Hyperlink(r, categories,front)
    r = categories.shapes[1].text_frame.paragraphs[i].add_run()
    r.text = " "

    # add a link to the vibby to music slides
    if music:
        notes = front.notes_slide.notes_text_frame.add_paragraph()
        notes.text = content[cat_index[i]+2]


    # make question slides themselves
    offset = 2+music
    qs_index = cat_index[i]+offset
    if i == len(cat_index)-1:
        qs_len = (len(content)-cat_index[i]-offset)//2
    else:
        qs_len = (cat_index[i+1]-cat_index[i]-offset)//2
    # music rounds only have question 10
    if music:
        qs_len = 1

    # create a slide with question number, and question
    for j in range(qs_len):
        q = prs.slides.add_slide(layout)
        p = q.placeholders[0].text_frame.add_paragraph()
        p.text = "Question " + str(j+1)
        # if music, question number is 10
        if music:
            p.text = "Question 10"
        p.font.size = Pt(44)
        # if picture, then make two slides; small picture below question number, and big picture
        if picture and j != qs_len - 1:
            filename = 'PhotoRounds/'+name+"/"+content[qs_index+2*j]+".png"
            q.shapes.add_picture(filename,Inches(1),Inches(2),
                                 height=Inches(prs.slide_height/914400-3))
            image_width = q.shapes[2].width/914400
            slide_width = prs.slide_width/914400
            if image_width > 0.8*slide_width:
                q.shapes[2].width = Inches(0.8*slide_width)
                image_width = 0.8*slide_width
            q.shapes[2].left = Inches((slide_width-image_width)/2)
            full_size = prs.slides.add_slide(layout)
            full_size.shapes.add_picture(filename,Inches(1),Inches(0.5),
                                         height = Inches(prs.slide_height/914400-1))
            image_width = full_size.shapes[2].width/914400
            if image_width > 0.9*slide_width:
                image_width = 0.9*slide_width
                full_size.shapes[2].width = Inches(0.9*slide_width)
            full_size.shapes[2].left = Inches((slide_width - image_width)/2)
        # otherwise just one slide with question number and body of question
        else: 
            ques = q.shapes[1].text_frame.paragraphs[0].add_run()
            ques.text = content[qs_index+2*j]
            ques.font.size = Pt(28)
        # if this is the last question, add a hyperlink back to Today's Categories
        if j == qs_len - 1:
            hyp = q.shapes[1].text_frame.add_paragraph().add_run()
            hyp.text = "Back to categories"
            Hyperlink(hyp, q, categories)
            

# make answer slides
answers = prs.slides.add_slide(layout)
answers_title = answers.placeholders[0].text_frame.add_paragraph()
answers_title.text = "Answers"
answers_title.font.size = Pt(44)

# answers for pictures and music are very different, so do them separate
for i in set(range(len(cat_index)))-set(picture_cats)-set(music_cats):
    # make a title slide for the category answers, and link to it from the Today's Categories
    front = prs.slides.add_slide(layout)
    category_slides.append(front)
    p = front.placeholders[0].text_frame.add_paragraph()
    p.text = content[cat_index[i]]
    p.font.size = Pt(44)
    desc = front.shapes[1].text_frame.paragraphs[0].add_run()
    desc.text = content[cat_index[i]+1]
    desc.font.size = Pt(28)
    
    r = categories.shapes[1].text_frame.paragraphs[i].add_run()
    r.text = "Answers"
    Hyperlink(r, categories,front)

    offset = 2+music
    qs_index = cat_index[i]+offset
    if i == len(cat_index)-1:
        qs_len = (len(content)-cat_index[i]-offset) // 2
    else:
        qs_len = (cat_index[i+1]-cat_index[i]-offset) // 2
    
    # reprint the question slide, then duplicate with a red answer below
    for j in range(qs_len):
        q = prs.slides.add_slide(layout)
        p = q.placeholders[0].text_frame.add_paragraph()
        p.text = "Question " + str(j+1)
        p.font.size = Pt(44)
        ques = q.shapes[1].text_frame.paragraphs[0].add_run()
        ques.text = content[qs_index+2*j]
        ques.font.size = Pt(28)
        q = prs.slides.add_slide(layout)
        p = q.placeholders[0].text_frame.add_paragraph()
        p.text = "Question " + str(j+1)
        p.font.size = Pt(44)
        ques = q.shapes[1].text_frame.paragraphs[0].add_run()
        ques.text = content[qs_index+2*j]
        ques.font.size = Pt(28)
        ans = q.shapes[1].text_frame.add_paragraph()
        ans.text = content[qs_index+2*j+1]
        ans.font.size = Pt(28)
        ans.font.color.rgb = RGBColor(255,0,0)

        # if last answer, hyperlink back to Today's Categories
        if j == qs_len - 1:
            hyp = q.shapes[1].text_frame.add_paragraph().add_run()
            hyp.text = "Back to categories"
            Hyperlink(hyp, q, categories)

# for picture answers, start with same title slide
for i in picture_cats:
    front = prs.slides.add_slide(layout)
    category_slides.append(front)
    p = front.placeholders[0].text_frame.add_paragraph()
    p.text = content[cat_index[i]]
    p.font.size = Pt(44)
    desc = front.shapes[1].text_frame.paragraphs[0].add_run()
    desc.text = content[cat_index[i]+1]
    desc.font.size = Pt(28)

    r = categories.shapes[1].text_frame.paragraphs[i].add_run()
    r.text = "Answers"
    Hyperlink(r, categories,front)

    
    qs_index = cat_index[i]+2

    # first slide is question picture, second slide is answer picture and text 
    if i == len(cat_index)-1:
        qs_len = (len(content)-cat_index[i]-offset) // 2
    else:
        qs_len = (cat_index[i+1]-cat_index[i]-offset) // 2
    for j in range(qs_len-1):
        q = prs.slides.add_slide(layout)
        p = q.placeholders[0].text_frame.add_paragraph()
        p.text = "Question " + str(j+1)
        p.font.size = Pt(44)

        filename = 'PhotoRounds/'+name+"/"+content[qs_index+2*j]
        print(filename)
        q.shapes.add_picture(filename+".png",Inches(1),Inches(2),
                             height=Inches(prs.slide_height/914400-3))
        image_width = q.shapes[2].width/914400
        slide_width = prs.slide_width/914400
        if image_width > 0.8*slide_width:
            q.shapes[2].width = Inches(0.8*slide_width)
            image_width = 0.8*slide_width
        q.shapes[2].left = Inches((slide_width-image_width)/2)

        ans = prs.slides.add_slide(layout)
        p = ans.placeholders[0].text_frame.add_paragraph()
        p.text = "Question " + str(j+1)
        p.font.size = Pt(44)

        ans.shapes.add_picture(filename+".jpg",Inches(1),Inches(2),
                             height=Inches(prs.slide_height/914400-3))
        image_width = ans.shapes[2].width/914400
        slide_width = prs.slide_width/914400
        if image_width > 0.8*slide_width:
            ans.shapes[2].width = Inches(0.8*slide_width)
            image_width = 0.8*slide_width
        ans.shapes[2].left = Inches((slide_width-image_width)/2)

        ans.shapes.add_textbox(Inches(1),Inches(prs.slide_height/914400-1),Inches(3),Inches(1))
        p = ans.shapes[3].text_frame.paragraphs[0].add_run()
        p.text = content[qs_index+2*j+1]
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255,0,0)

    # last question isn't picture; add question on one slide, answer on next
    k = qs_len-1
    q = prs.slides.add_slide(layout)
    p = q.placeholders[0].text_frame.add_paragraph()
    p.text = "Question " + str(qs_len)
    p.font.size = Pt(44)
    ques = q.shapes[1].text_frame.paragraphs[0].add_run()
    ques.text = content[qs_index+2*k]
    ques.font.size = Pt(28)
    q = prs.slides.add_slide(layout)
    p = q.placeholders[0].text_frame.add_paragraph()
    p.text = "Question " + str(k+1)
    p.font.size = Pt(44)
    ques = q.shapes[1].text_frame.paragraphs[0].add_run()
    ques.text = content[qs_index+2*k]
    ques.font.size = Pt(28)
    ans = q.shapes[1].text_frame.add_paragraph()
    ans.text = content[qs_index+2*k+1]
    ans.font.size = Pt(28)
    ans.font.color.rgb = RGBColor(255,0,0)

    hyp = q.shapes[1].text_frame.add_paragraph().add_run()
    hyp.text = "Back to categories"
    Hyperlink(hyp, q, categories)  

# answers to music
for i in music_cats:
    front = prs.slides.add_slide(layout)
    category_slides.append(front)
    p = front.placeholders[0].text_frame.add_paragraph()
    p.text = content[cat_index[i]]
    p.font.size = Pt(44)
    desc = front.shapes[1].text_frame.paragraphs[0].add_run()
    desc.text = content[cat_index[i]+1]
    desc.font.size = Pt(28)

    r = categories.shapes[1].text_frame.paragraphs[i].add_run()
    r.text = "Answers"
    Hyperlink(r, categories,front)


    qs_index = cat_index[i]+5

    if i == len(cat_index)-1:
        qs_len = (len(content)-cat_index[i]-5)
    else:
        qs_len = (cat_index[i+1]-cat_index[i]-5)

    q = prs.slides.add_slide(layout)
    p = q.placeholders[0].text_frame.add_paragraph()
    p.text = "Music Answers"
    p.font.size = Pt(44)

    #answers for q 1 through 9
    for j in range(qs_len):
        q = prs.slides.add_slide(layout)
        p = q.placeholders[0].text_frame.add_paragraph()
        p.text = "Music Answers"
        p.font.size = Pt(44)
        ans = q.shapes[1].text_frame.paragraphs[0].add_run()
        ans.text = content[qs_index]
        ans.font.size = Pt(28)
        ans.font.color.rgb=RGBColor(255,0,0)
        for k in range(1,j+1):
            ans = q.shapes[1].text_frame.add_paragraph()
            ans.text = content[qs_index+k]
            ans.font.size = Pt(28)
            ans.font.color.rgb = RGBColor(255,0,0)

    # answer to question 10
    q = prs.slides.add_slide(layout)
    p = q.placeholders[0].text_frame.add_paragraph()
    p.text = "Question 10"
    p.font.size = Pt(44)
    ques = q.placeholders[1].text_frame.paragraphs[0].add_run()
    ques.text = content[qs_index-2]
    ques.font.size = Pt(28)
    q = prs.slides.add_slide(layout)
    p =q.placeholders[0].text_frame.add_paragraph()
    p.text = "Question 10"
    p.font.size= Pt(44)
    ques = q.placeholders[1].text_frame.paragraphs[0].add_run()
    ques.text =content[qs_index-2]
    ques.font.size = Pt(28)
    ans = q.shapes[1].text_frame.add_paragraph()
    ans.text = content[qs_index-1]
    ans.font.size = Pt(28)
    ans.font.color.rgb = RGBColor(255,0,0)


        

for slide in prs.slides:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(141,167,222)

for slide in category_slides:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(86, 182, 69)

answers.background.fill.solid()
answers.background.fill.fore_color.rgb = RGBColor(255,0,0)






prs.save(name+".pptx")
